from subprocess import call
from pptx import Presentation
from unbabel.api import UnbabelApi

api = UnbabelApi(username='andre.nascimento.melo',api_key='d7863bba9180c360e76cdc1bccdbae9705427c08',sandbox=True)

def get_presentation_text(presentation_file_path):
	presentation_file = open(presentation_file_path)
	presentation = Presentation(presentation_file)
	presentation_file.close()

	text_runs = []

	for slide in presentation.slides:
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			for paragraph in shape.text_frame.paragraphs:
				for run in paragraph.runs:
					# print(run.text)
					if len(str(run.text.encode('utf-8')).strip())>0:
						text_runs.append(run.text)
	                	
	map(lambda x: x.encode('utf-8'), text_runs)
	return text_runs

def translate_list_of_text(content, target_lang, source_lang):
	translations=[]
	
	for text_to_translate in content:
		t=api.post_translations(text=text_to_translate, target_language = target_lang, source_language = source_lang)
		translations.append(t)

	print "returning pending translations with len " + str(len(translations))
	return list(translations)

def translate_power_point(room_id, translations, target_lang):
	presentation_file_path = "./uploads/"+room_id+"/"+room_id+".pptx"
	presentation_file = open(presentation_file_path)
	presentation = Presentation(presentation_file)
	presentation_file.close()

	for slide in presentation.slides:
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			for paragraph in shape.text_frame.paragraphs:
				for run in paragraph.runs:
					if len(str(run.text.encode('utf-8')).strip())>0:
						run.text=translations.pop(0)

	call(["mkdir","./uploads/"+room_id+'_'+target_lang])
	presentation.save("./uploads/"+room_id+'_'+target_lang+'/'+room_id+'_'+target_lang + '.pptx')
	call(["unoconv", "-f", "html", "./uploads/"+room_id+'_'+target_lang+'/'+room_id+'_'+target_lang+'.pptx'])